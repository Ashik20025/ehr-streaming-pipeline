from __future__ import annotations

import json
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ARTIFACTS = ROOT / "artifacts"
ASSETS = ARTIFACTS / "presentation_assets"

TITLE = "Real-Time Streaming EHR Ingestion and Validation Pipeline"
AUTHOR = "Hanuman Ashik Akshintala"
SUBTITLE = 'Inspired by "Real-Time Data Ingestion Pipelines for Streaming EHR Systems"'

GREEN = RGBColor(0x00, 0x79, 0x6B)
RED = RGBColor(0xBD, 0x3F, 0x53)
GOLD = RGBColor(0xAD, 0x7C, 0x00)
DARK = RGBColor(0x17, 0x21, 0x1F)
MUTED = RGBColor(0x51, 0x62, 0x5D)
LIGHT = RGBColor(0xF4, 0xF8, 0xF6)
BORDER = RGBColor(0xD8, 0xE3, 0xDF)
SLATE = RGBColor(0x43, 0x5B, 0x57)


def main() -> None:
    latest = json.loads((ARTIFACTS / "distributed_summary.json").read_text(encoding="utf-8"))
    experiments = json.loads((ARTIFACTS / "experiments" / "experiment_summary.json").read_text(encoding="utf-8"))

    output = ARTIFACTS / "EECS247_Final_Presentation_Hanuman_Akshintala.pptx"
    deck = build_deck(latest, experiments)
    deck.save(output)

    downloads_copy = Path("/Users/hanumanashikakshintala/Downloads/EECS247_Final_Presentation_Hanuman_Akshintala.pptx")
    downloads_copy.write_bytes(output.read_bytes())
    print(f"wrote {output}")
    print(f"copied {downloads_copy}")


def build_deck(latest: dict, experiments: list[dict]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_problem_slide(prs)
    add_architecture_slide(prs)
    add_implementation_slide(prs)
    add_validation_slide(prs)
    add_experiment_setup_slide(prs)
    add_results_highlights_slide(prs, latest, experiments)
    add_experiment_summary_slide(prs)
    add_dashboard_slide(prs, latest)
    add_runtime_and_output_slide(prs)
    add_discussion_slide(prs)
    add_demo_slide(prs)
    add_closing_slide(prs)
    return prs


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)

    textbox(slide, 0.8, 0.8, 11.7, 1.4, TITLE, 28, True, DARK)
    textbox(slide, 0.9, 2.2, 8.8, 0.7, AUTHOR, 20, True, GREEN)
    textbox(slide, 0.9, 2.75, 10.4, 0.7, "EECS 247 Course Project Presentation", 18, False, DARK)
    textbox(slide, 0.9, 3.3, 11.2, 0.8, SUBTITLE, 18, False, MUTED)
    textbox(slide, 0.9, 3.9, 10.0, 0.6, f"April 29, 2026", 16, False, MUTED)

    add_band(slide, 0.8, 5.1, 11.7, 1.4, GREEN)
    bullets(
        slide,
        1.1,
        5.35,
        11.0,
        0.95,
        [
            "Local streaming pipeline using Kafka, Flink, Redis, Cassandra, and a React/Material UI dashboard",
            "Synthetic EHR events with validation, deduplication, quarantine, and performance reporting",
        ],
        18,
        LIGHT,
    )


def add_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Problem and Motivation")
    bullets(
        slide,
        0.9,
        1.45,
        6.0,
        4.0,
        [
            "Modern healthcare systems increasingly receive EHR data as streams rather than offline batches.",
            "Bad records can silently corrupt downstream analytics and decision-support systems if there is no validation path.",
            "Streaming pipelines need duplicate detection, late-event handling, and clear quarantine routing.",
            "The project goal was to implement a local version of a paper-inspired real-time EHR ingestion system.",
        ],
        20,
        DARK,
    )
    add_stat_card(slide, 7.4, 1.7, 2.0, 1.4, "4", "experiment runs", GREEN)
    add_stat_card(slide, 9.7, 1.7, 2.0, 1.4, "30,000", "events tested", GOLD)
    add_stat_card(slide, 7.4, 3.45, 2.0, 1.4, "94.25%", "accepted", SLATE)
    add_stat_card(slide, 9.7, 3.45, 2.0, 1.4, "5.75%", "quarantined", RED)
    textbox(slide, 7.4, 5.2, 4.3, 1.2, "All runs use synthetic EHR events and execute locally without cloud services or real patient data.", 17, False, MUTED)


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "System Architecture")

    labels = [
        ("Synthetic EHR\nProducer", GREEN),
        ("Kafka", GREEN),
        ("Flink", GREEN),
        ("Redis", GOLD),
        ("Cassandra", GREEN),
        ("Dashboard /\nReports", SLATE),
    ]
    left = 0.8
    top = 2.1
    width = 1.8
    height = 1.0
    gap = 0.35
    boxes = []
    for label, color in labels:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = LIGHT
        p.alignment = PP_ALIGN.CENTER
        boxes.append((left, top, width, height))
        left += width + gap

    for index in range(len(boxes) - 1):
        x1 = boxes[index][0] + boxes[index][2]
        y = boxes[index][1] + boxes[index][3] / 2
        x2 = boxes[index + 1][0]
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
        connector.line.color.rgb = MUTED
        connector.line.width = Pt(2.5)

    bullets(
        slide,
        0.95,
        4.1,
        11.3,
        2.4,
        [
            "Kafka carries the raw synthetic EHR stream and separates valid, quarantine, and alert topics.",
            "Flink performs validation, keyed-state duplicate detection, event-time processing, enrichment, and alert logic.",
            "Redis stores low-latency patient metadata. Cassandra stores accepted records, quarantined records, and alerts.",
        ],
        18,
        DARK,
    )


def add_implementation_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Implemented Components")

    cards = [
        ("Synthetic Generator", "Vitals, labs, medications, admissions; injects duplicates, late events, and invalid records.", GREEN),
        ("Flink Validation", "Schema checks, value-range checks, timestamp parsing, keyed duplicate detection, watermarks.", SLATE),
        ("Quarantine Path", "Bad records are stored separately with explicit reasons instead of being silently dropped.", RED),
        ("Serving Layer", "Redis for enrichment metadata; Cassandra for valid events, quarantine, and alert tables.", GOLD),
    ]
    positions = [(0.8, 1.7), (6.85, 1.7), (0.8, 4.0), (6.85, 4.0)]
    for (label, body, color), (x, y) in zip(cards, positions):
        add_info_card(slide, x, y, 5.6, 1.8, label, body, color)


def add_validation_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Validation and Stream Logic")
    bullets(
        slide,
        0.9,
        1.55,
        6.3,
        4.5,
        [
            "Required-field validation: event ID, patient ID, event type, event time, schema version",
            "Clinical value validation: vital and lab out-of-range detection",
            "Duplicate detection: keyed state over event IDs in Flink",
            "Late and out-of-order handling: event-time watermarks with allowed lateness window",
            "PHI-like field rejection: direct patient-name fields are quarantined",
            "Possible-sepsis alert generation from the valid-event stream",
        ],
        18,
        DARK,
    )

    add_stat_card(slide, 8.0, 1.9, 2.1, 1.3, "175", "top duplicate count", RED)
    add_stat_card(slide, 10.4, 1.9, 2.1, 1.3, "84", "direct PHI rejects", RED)
    add_stat_card(slide, 8.0, 3.55, 2.1, 1.3, "83", "missing patient ID", RED)
    add_stat_card(slide, 10.4, 3.55, 2.1, 1.3, "80", "unknown event type", RED)
    textbox(slide, 7.95, 5.35, 4.6, 1.0, "Top quarantine reasons come from the 10,000-event final run with unlimited producer rate.", 15, False, MUTED)


def add_experiment_setup_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Experiment Setup")
    bullets(
        slide,
        0.95,
        1.55,
        5.4,
        3.4,
        [
            "Synthetic patients: 5,000",
            "Allowed lateness window: 300 seconds",
            "Invalid record injection: 4%",
            "Duplicate injection: 2%",
            "Out-of-order injection: 12%",
            "Late-event injection: 3%",
        ],
        18,
        DARK,
    )

    add_table_block(
        slide,
        6.3,
        1.75,
        6.0,
        3.5,
        [
            ["Run", "Events", "Rate limit"],
            ["events_5000_rate_1000", "5,000", "1,000/s"],
            ["events_5000_rate_0", "5,000", "as fast as possible"],
            ["events_10000_rate_1000", "10,000", "1,000/s"],
            ["events_10000_rate_0", "10,000", "as fast as possible"],
        ],
    )
    textbox(slide, 0.95, 5.45, 11.0, 0.8, "The goal of the matrix was to compare controlled-rate ingestion against maximum-speed producer throughput and observe the effect on latency.", 16, False, MUTED)


def add_results_highlights_slide(prs: Presentation, latest: dict, experiments: list[dict]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Results Highlights")

    best_throughput = max(experiments, key=lambda row: row["producer_events_per_sec"])
    best_latency = min(experiments, key=lambda row: row["p95_latency_ms"])
    latest_totals = latest["totals"]
    latest_latency = latest["latency_ms"]

    add_stat_card(slide, 0.9, 1.7, 2.7, 1.5, f'{best_throughput["producer_events_per_sec"]:,.0f}/s', "best throughput", GREEN)
    add_stat_card(slide, 3.95, 1.7, 2.7, 1.5, f'{best_latency["p95_latency_ms"]:,.0f} ms', "lowest p95 latency", GOLD)
    add_stat_card(slide, 7.0, 1.7, 2.7, 1.5, f'{latest_totals["accepted_percent"]:.2f}%', "accepted records", SLATE)
    add_stat_card(slide, 10.05, 1.7, 2.3, 1.5, f'{latest_totals["alerts"]}', "alerts", RED)

    bullets(
        slide,
        0.95,
        3.9,
        11.3,
        2.3,
        [
            f'Unlimited producer rate increased throughput to {best_throughput["producer_events_per_sec"]:,.0f} events/sec but also increased p95 latency.',
            f'Controlled-rate ingestion at 1,000 events/sec achieved the lowest p95 latency of {best_latency["p95_latency_ms"]:,.0f} ms.',
            f'The latest 10,000-event run accepted {int(latest_totals["valid"]):,} records and quarantined {int(latest_totals["quarantine"]):,}.',
            f'Latest-run average latency was {latest_latency["avg_ms"]:,.0f} ms and p95 latency was {latest_latency["p95_ms"]:,.0f} ms.',
        ],
        18,
        DARK,
    )


def add_experiment_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Experiment Matrix Summary")
    textbox(slide, 0.95, 1.1, 11.2, 0.5, "Final table aggregating the four workload configurations.", 16, False, MUTED)
    add_picture_fit(slide, ASSETS / "experiment_summary.png", 0.7, 1.7, 12.0, 5.1)


def add_dashboard_slide(prs: Presentation, latest: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "React Dashboard")
    textbox(slide, 0.95, 1.1, 11.2, 0.5, "Live React + Material UI dashboard served at localhost:5173.", 16, False, MUTED)
    add_picture_fit(slide, ASSETS / "react_dashboard.png", 0.6, 1.6, 12.1, 4.9)
    latest_totals = latest["totals"]
    latest_latency = latest["latency_ms"]
    bullets(
        slide,
        0.95,
        6.0,
        12.0,
        1.0,
        [
            f'The live dashboard shows {int(latest_totals["valid"]):,} valid records, {int(latest_totals["quarantine"]):,} quarantined records, and {int(latest_totals["alerts"]):,} alert(s).',
            f'Latest-run p95 latency is {latest_latency["p95_ms"]:,.0f} ms, and throughput is derived directly from the generated pipeline summary.',
        ],
        15,
        DARK,
    )


def add_runtime_and_output_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Runtime and Generated Output")
    add_picture_fit(slide, ASSETS / "flink_dashboard.png", 0.7, 1.4, 5.9, 3.75)
    add_picture_fit(slide, ASSETS / "distributed_report.png", 6.8, 1.4, 5.85, 3.75)
    textbox(slide, 0.8, 5.35, 5.7, 0.4, "Flink dashboard confirming the streaming job is running.", 15, True, DARK)
    textbox(slide, 6.9, 5.35, 5.5, 0.4, "Generated report summarizing final output from Cassandra summaries.", 15, True, DARK)
    bullets(
        slide,
        0.9,
        5.85,
        11.7,
        0.75,
        [
            "The Flink UI shows the live streaming job. The report view is a stable output for final analysis and presentation backup.",
            "Quarantined records are expected because the generator intentionally injects duplicates, malformed timestamps, PHI-like fields, and range violations.",
        ],
        14,
        DARK,
    )


def add_discussion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Discussion, Limitations, and Future Work")
    bullets(
        slide,
        0.95,
        1.55,
        5.7,
        4.9,
        [
            "Why Redis: fast metadata enrichment during stream processing.",
            "Why Cassandra: scalable write-heavy serving store for valid events, quarantine, and alerts.",
            "Why latency rises at unlimited rate: the producer can outrun downstream storage and processing.",
            "Limitations: synthetic data, single-node local deployment, and simple alert logic.",
            "Future work: Synthea/FHIR input, restart/failure experiments, and more advanced alert rules.",
        ],
        18,
        DARK,
    )
    add_stat_card(slide, 7.8, 1.9, 2.2, 1.4, "5,000", "patients", GREEN)
    add_stat_card(slide, 10.2, 1.9, 2.2, 1.4, "300s", "lateness window", GOLD)
    add_stat_card(slide, 7.8, 3.7, 2.2, 1.4, "4%", "invalid injection", RED)
    add_stat_card(slide, 10.2, 3.7, 2.2, 1.4, "12%", "out-of-order", SLATE)


def add_demo_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    title(slide, "Presentation Demo Plan")
    bullets(
        slide,
        0.95,
        1.55,
        11.2,
        3.8,
        [
            "Start with the architecture slide and explain the Kafka -> Flink -> Redis/Cassandra data path.",
            "Show the experiment summary table to compare 5,000 vs 10,000 events and rate-limited vs unlimited ingestion.",
            "Open the generated output report to explain valid records, quarantine reasons, alerts, and event-type counts.",
            "If Docker Desktop is available, open the live dashboard at http://localhost:5173 and Flink at http://localhost:8081.",
        ],
        18,
        DARK,
    )
    add_code_box(
        slide,
        1.0,
        4.95,
        11.0,
        1.3,
        [
            'cd "/Users/hanumanashikakshintala/Documents/New project/ehr_stream_pipeline"',
            "./scripts/start_presentation_demo.sh",
        ],
    )


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, LIGHT)
    textbox(slide, 1.0, 1.5, 11.2, 1.1, "Takeaway", 28, True, DARK)
    textbox(
        slide,
        1.0,
        2.4,
        11.0,
        2.0,
        "The project demonstrates that a paper-inspired streaming EHR ingestion pipeline can be implemented locally with real Kafka, Flink, Redis, and Cassandra components while producing measurable throughput, latency, and quarantine behavior.",
        22,
        False,
        DARK,
    )
    textbox(slide, 1.0, 5.1, 11.0, 0.8, "Questions?", 26, True, GREEN)


def background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title(slide, text: str) -> None:
    textbox(slide, 0.8, 0.45, 11.8, 0.6, text, 24, True, DARK)


def textbox(slide, x: float, y: float, w: float, h: float, text: str, size: int, bold: bool, color: RGBColor):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def bullets(slide, x: float, y: float, w: float, h: float, items: list[str], size: int, color: RGBColor) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_band(slide, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_stat_card(slide, x: float, y: float, w: float, h: float, value: str, label: str, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = BORDER
    textbox(slide, x + 0.15, y + 0.12, w - 0.3, 0.35, label, 14, False, MUTED)
    textbox(slide, x + 0.15, y + 0.48, w - 0.3, 0.55, value, 22, True, DARK)
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.color.rgb = accent


def add_info_card(slide, x: float, y: float, w: float, h: float, heading: str, body: str, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = BORDER
    textbox(slide, x + 0.18, y + 0.15, w - 0.36, 0.35, heading, 18, True, accent)
    textbox(slide, x + 0.18, y + 0.58, w - 0.36, 1.05, body, 15, False, DARK)


def add_table_block(slide, x: float, y: float, w: float, h: float, rows: list[list[str]]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = BORDER
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x + 0.12), Inches(y + 0.12), Inches(w - 0.24), Inches(h - 0.24)).table
    col_widths = [2.65, 1.0, 1.6]
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)
    row_height = (h - 0.24) / len(rows)
    for row in table.rows:
        row.height = Inches(row_height)
    for r, row_values in enumerate(rows):
        for c, value in enumerate(row_values):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14 if r == 0 else 13)
                p.font.bold = r == 0
                p.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xED, 0xF5, 0xF2) if r == 0 else RGBColor(0xFF, 0xFF, 0xFF)


def add_code_box(slide, x: float, y: float, w: float, h: float, lines: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.color.rgb = DARK
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.36), Inches(h - 0.24))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Menlo"
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT


def add_picture_fit(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


if __name__ == "__main__":
    main()
